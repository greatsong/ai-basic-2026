#!/usr/bin/env python3
"""
AI 기초 과정중심평가 문항집 PPTX 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === CONSTANTS ===
PRIMARY_BLUE = RGBColor(0x3B, 0x82, 0xF6)
DARK_BLUE = RGBColor(0x1E, 0x40, 0xAF)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
VERY_LIGHT_BLUE = RGBColor(0xEF, 0xF6, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x6B, 0x72, 0x80)
DARK_GRAY = RGBColor(0x37, 0x41, 0x51)
GREEN = RGBColor(0x05, 0x96, 0x69)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
RED = RGBColor(0xDC, 0x26, 0x26)

IMG_DIR = "/Users/greatsong/greatsong-project/ai-basic-2026/assessment-images"
OUT_PATH = "/Users/greatsong/greatsong-project/ai-basic-2026/assessment-slides.pptx"

FONT_NAME = "Malgun Gothic"
FONT_NAME_EN = "Arial"

# Slide dimensions: 16:9
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# === HELPER FUNCTIONS ===

def add_bg_rect(slide, color=VERY_LIGHT_BLUE):
    """Add a full-slide background rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_top_bar(slide, color=PRIMARY_BLUE, height=Inches(0.08)):
    """Add a thin top accent bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bottom_bar(slide, text="AI 기초 과정중심평가 문항집", color=DARK_BLUE):
    """Add a bottom bar with text."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SLIDE_HEIGHT - Inches(0.45), SLIDE_WIDTH, Inches(0.45)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER


def set_text(text_frame, text, font_size=18, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    """Set text in a text frame."""
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment


def add_paragraph(text_frame, text, font_size=16, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME, space_before=Pt(4), space_after=Pt(4)):
    """Add a paragraph to a text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    if space_after:
        p.space_after = space_after
    return p


def add_textbox(slide, left, top, width, height, text="", font_size=18, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    """Add a textbox to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    set_text(tf, text, font_size, bold, color, alignment, font_name)
    return txBox


def add_image_safe(slide, img_name, left, top, width=None, height=None):
    """Add an image if it exists. Returns the shape or None."""
    img_path = os.path.join(IMG_DIR, img_name)
    if os.path.exists(img_path):
        kwargs = {}
        if width:
            kwargs['width'] = width
        if height:
            kwargs['height'] = height
        if not width and not height:
            kwargs['height'] = Inches(2.5)
        try:
            pic = slide.shapes.add_picture(img_path, left, top, **kwargs)
            return pic
        except Exception as e:
            print(f"Warning: Could not add image {img_name}: {e}")
            return None
    else:
        print(f"Warning: Image not found: {img_path}")
        return None


def add_tag_badge(slide, tag_text, left, top, color=PRIMARY_BLUE):
    """Add a colored tag badge."""
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.5), Inches(0.35)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = False
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = tag_text
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME
    p.font.bold = True
    return badge


THEME_COLORS = {
    "CT": RGBColor(0xDC, 0x26, 0x26),  # red
    "CO": PRIMARY_BLUE,
    "AL": GREEN,
    "MA": ORANGE,
}


def add_theme_badge(slide, theme, lesson, left, top):
    """Add theme and lesson badges."""
    color = THEME_COLORS.get(theme, GRAY)
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(0.8), Inches(0.32)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = theme
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME_EN
    p.font.bold = True

    lesson_badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.9), top, Inches(1.2), Inches(0.32)
    )
    lesson_badge.fill.solid()
    lesson_badge.fill.fore_color.rgb = GRAY
    lesson_badge.line.fill.background()
    tf2 = lesson_badge.text_frame
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = tf2.paragraphs[0]
    p2.text = lesson
    p2.font.size = Pt(11)
    p2.font.color.rgb = WHITE
    p2.font.name = FONT_NAME
    p2.font.bold = False


def create_title_slide(title, subtitle):
    """Create a styled title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = DARK_BLUE
    bg_shape.line.fill.background()

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.2), Inches(2), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_BLUE
    line.line.fill.background()

    # Title
    add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.6),
                title, font_size=40, bold=True, color=WHITE, font_name=FONT_NAME)

    # Subtitle
    add_textbox(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(1),
                subtitle, font_size=20, color=RGBColor(0xBF, 0xDB, 0xFE), font_name=FONT_NAME)

    return slide


def create_part_title_slide(part_code, title, point_info=""):
    """Create a part separator slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = PRIMARY_BLUE
    bg_shape.line.fill.background()

    # Part code (big)
    add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
                part_code, font_size=60, bold=True, color=WHITE, font_name=FONT_NAME_EN)

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.0), Inches(2.5), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = WHITE
    line.line.fill.background()

    # Title
    add_textbox(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(1),
                title, font_size=28, bold=True, color=WHITE, font_name=FONT_NAME)

    # Point info
    if point_info:
        add_textbox(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.6),
                    point_info, font_size=18, color=RGBColor(0xBF, 0xDB, 0xFE), font_name=FONT_NAME)

    return slide


def create_question_slide(q_num, theme, lesson, question_text, img_name=None, extra_text=None):
    """Create a standard question slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)

    # Question number
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(3), Inches(0.6),
                q_num, font_size=28, bold=True, color=DARK_BLUE, font_name=FONT_NAME_EN)

    # Theme and lesson badges
    add_theme_badge(slide, theme, lesson, Inches(0.8), Inches(1.0))

    # Question text area
    text_left = Inches(0.8)
    text_top = Inches(1.6)
    text_width = Inches(7) if img_name else Inches(11.5)
    text_height = Inches(4.8)

    txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Split question text into lines
    lines = question_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK
        p.font.name = FONT_NAME
        p.space_after = Pt(6)

    if extra_text:
        p = tf.add_paragraph()
        p.text = ""
        p.space_after = Pt(8)
        p = tf.add_paragraph()
        p.text = extra_text
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.font.name = FONT_NAME
        p.font.italic = True

    # Image
    if img_name:
        img_left = Inches(8.3)
        img_top = Inches(1.2)
        add_image_safe(slide, img_name, img_left, img_top, height=Inches(4))

    return slide


def create_question_slide_with_two_images(q_num, theme, lesson, question_text, img1, img2):
    """Create a question slide with two images side by side."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(3), Inches(0.6),
                q_num, font_size=28, bold=True, color=DARK_BLUE, font_name=FONT_NAME_EN)

    add_theme_badge(slide, theme, lesson, Inches(0.8), Inches(1.0))

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(6.5), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = question_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = BLACK
        p.font.name = FONT_NAME
        p.space_after = Pt(6)

    add_image_safe(slide, img1, Inches(7.8), Inches(1.2), height=Inches(2.5))
    add_image_safe(slide, img2, Inches(7.8), Inches(4.0), height=Inches(2.5))

    return slide


def create_calculation_slide(q_num, theme, lesson, problem_text, solution_text, img_name=None):
    """Create a calculation problem slide with solution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, WHITE)
    add_top_bar(slide, color=ORANGE)
    add_bottom_bar(slide)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(3), Inches(0.6),
                q_num, font_size=28, bold=True, color=ORANGE, font_name=FONT_NAME_EN)

    add_theme_badge(slide, theme, lesson, Inches(0.8), Inches(1.0))

    # Problem
    text_width = Inches(6.5) if img_name else Inches(11.5)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), text_width, Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = problem_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = BLACK
        p.font.name = FONT_NAME
        p.space_after = Pt(4)

    # Solution box
    sol_top = Inches(4.2)
    sol_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), sol_top, text_width, Inches(2.5)
    )
    sol_box.fill.solid()
    sol_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF7, 0xED)
    sol_box.line.color.rgb = ORANGE
    sol_box.line.width = Pt(1)

    tf2 = sol_box.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.2)
    tf2.margin_top = Inches(0.15)

    p = tf2.paragraphs[0]
    p.text = "풀이"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.font.name = FONT_NAME

    sol_lines = solution_text.strip().split('\n')
    for line in sol_lines:
        p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.font.name = FONT_NAME
        p.space_after = Pt(3)

    if img_name:
        add_image_safe(slide, img_name, Inches(8), Inches(1.2), height=Inches(4))

    return slide


def create_essay_slide(q_num, theme, lesson, question_text, keywords, img_name=None, img_name2=None):
    """Create an essay question slide with keywords."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, WHITE)
    add_top_bar(slide, color=DARK_BLUE)
    add_bottom_bar(slide)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(3), Inches(0.6),
                q_num, font_size=28, bold=True, color=DARK_BLUE, font_name=FONT_NAME_EN)

    add_theme_badge(slide, theme, lesson, Inches(0.8), Inches(1.0))

    text_width = Inches(7) if img_name else Inches(11.5)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), text_width, Inches(3.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = question_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK
        p.font.name = FONT_NAME
        p.space_after = Pt(6)

    # Keywords box
    kw_top = Inches(5.0)
    kw_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), kw_top, text_width, Inches(0.9)
    )
    kw_box.fill.solid()
    kw_box.fill.fore_color.rgb = LIGHT_BLUE
    kw_box.line.fill.background()

    tf_kw = kw_box.text_frame
    tf_kw.word_wrap = True
    tf_kw.margin_left = Inches(0.2)
    tf_kw.margin_top = Inches(0.1)
    p = tf_kw.paragraphs[0]
    p.text = f"핵심 키워드: {keywords}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = FONT_NAME

    if img_name:
        add_image_safe(slide, img_name, Inches(8.3), Inches(1.2), height=Inches(2.8))
    if img_name2:
        add_image_safe(slide, img_name2, Inches(8.3), Inches(4.2), height=Inches(2.5))

    return slide


# ============================================
# SLIDE CREATION
# ============================================

print("Creating slides...")

# --- 1. TITLE SLIDE ---
create_title_slide(
    "AI 기초 과정중심평가 문항집",
    "고등학교 2학년  |  인공지능 기초  |  2026"
)

# --- 2. OVERVIEW SLIDE ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_top_bar(slide)
add_bottom_bar(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
            "평가 개요 — 4대 핵심 평가 테마", font_size=28, bold=True, color=DARK_BLUE)

# Overview table as shapes
themes_data = [
    ("CT", "비판적 사고력", "AI의 가능성과 한계를 근거 기반으로 판단", RGBColor(0xDC, 0x26, 0x26)),
    ("CO", "개념 이해", "AI 핵심 구성 요소의 정의와 역할을 정확히 설명", PRIMARY_BLUE),
    ("AL", "알고리즘 이해", "AI 학습 알고리즘의 단계와 흐름을 절차적으로 설명", GREEN),
    ("MA", "수학적 이해", "AI 원리의 수학적 기초를 이해하고 직접 계산에 적용", ORANGE),
]

for i, (code, name, desc, color) in enumerate(themes_data):
    y = Inches(1.5) + Inches(1.3) * i
    # Code badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y, Inches(1.2), Inches(0.9))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME_EN

    # Name
    add_textbox(slide, Inches(2.5), y, Inches(3), Inches(0.5),
                name, font_size=20, bold=True, color=DARK_GRAY)
    # Description
    add_textbox(slide, Inches(2.5), y + Inches(0.45), Inches(9), Inches(0.5),
                desc, font_size=14, color=GRAY)

# Parts summary
parts_data = [
    ("A", "빈칸/단답", "10문항", "20점"),
    ("B", "계산형", "9문항", "계산"),
    ("C", "시각/분석형", "5문항", "16점"),
    ("D", "서술형(개별)", "6문항", "24점"),
    ("E", "서술형(통합)", "2문항", "10점"),
    ("F", "자기 성찰", "2문항", "Google Forms"),
]

parts_top = Inches(1.3)
for i, (code, typ, cnt, pts) in enumerate(parts_data):
    x = Inches(7.5)
    y = parts_top + Inches(0.85) * i
    add_textbox(slide, x, y, Inches(0.8), Inches(0.35),
                f"Part {code}", font_size=12, bold=True, color=DARK_BLUE, font_name=FONT_NAME_EN)
    add_textbox(slide, x + Inches(0.9), y, Inches(2), Inches(0.35),
                f"{typ} ({cnt})", font_size=12, color=DARK_GRAY)
    add_textbox(slide, x + Inches(3.2), y, Inches(1.5), Inches(0.35),
                pts, font_size=12, bold=True, color=GRAY)


# --- 3. PART A TITLE ---
create_part_title_slide("Part A", "기본 개념 (빈칸/단답)", "10문항 x 2점 = 20점")

# --- 4. A-1 through A-10 ---
# A-1
create_question_slide(
    "A-1", "CO", "1차시",
    "위 사진의 인물은 1950년 \"기계가 생각할 수 있는가?\"라는\n질문을 던진 영국의 수학자이다.\n\n이 사람의 이름은 (          )이다.",
    img_name="alan_turing.jpg"
)

# A-2
create_question_slide(
    "A-2", "CO", "2차시",
    "위 사진의 인물은 매컬럭과 핏츠(McCulloch & Pitts)의\n인공 뉴런 아이디어에 학습 규칙을 더하여\n새로운 기계를 만들었다.\n\n이 기계의 이름은 (          )이다.",
    img_name="rosenblatt.jpg"
)

# A-3
create_question_slide(
    "A-3", "CO", "3차시",
    "위 그래프 (다)처럼 출력이 0에서 1 사이의\nS자 곡선인 활성화함수의 이름은 (          )이다.",
    img_name="activation_functions_unlabeled.png"
)

# A-4
create_question_slide(
    "A-4", "CO", "3차시",
    "위 그래프 (가)처럼 입력이 0 이하이면 0,\n양수이면 그대로 출력하는 활성화함수의\n이름은 (          )이다.",
    img_name="activation_functions_unlabeled.png"
)

# A-5
create_question_slide(
    "A-5", "CO", "4차시",
    "위 그래프처럼, AI가 학습하면서 점점 줄어드는 값이 있다.\n예측값과 실제값의 차이를 하나의 숫자로 측정하는\n이 함수를 (          )이라 한다.",
    img_name="loss_decreasing.jpg"
)

# A-6
create_question_slide(
    "A-6", "CO", "5차시",
    "위 그래프에서, 한 번에 이동하는 폭이 너무 크면\n발산하고 너무 작으면 느리다.\n이 이동 폭을 결정하는 값을 (          )이라 한다.",
    img_name="gradient_descent_comparison.png"
)

# A-7
create_question_slide(
    "A-7", "AL", "2차시",
    "퍼셉트론의 학습 규칙은\n\"예측이 틀리면 가중치를 (          )하고,\n맞으면 (          )한다\"이다.",
    img_name="perceptron_diagram.png"
)

# A-8
create_question_slide(
    "A-8", "AL", "5차시",
    "경사하강법의 가중치 업데이트 공식:\n\nw_new = w_old - (          ) x (          )",
    img_name="gradient_descent_comparison.png"
)

# A-9
create_question_slide(
    "A-9", "CO", "7차시",
    "\"unhappiness\"를 \"un\", \"happi\", \"ness\"처럼\n의미 단위로 분절하는 위 그림의\n토큰화 방식을 (          )이라 한다.",
    img_name="bpe_tokenization.jpg"
)

# A-10
create_question_slide(
    "A-10", "CO", "8차시",
    "위 그림처럼, AI가 사실이 아닌 내용을\n자신 있게 생성하는 현상을 (          )이라 한다.",
    img_name="ai_hallucination.jpg"
)

# --- 5. PART B TITLE ---
create_part_title_slide("Part B", "수학적 이해 (계산형)", "9문항 | 풀이 과정 60% + 최종 답 40%")

# --- 6. B-1 through B-9 ---
# B-1
create_calculation_slide(
    "B-1", "MA", "2차시",
    "AND 게이트의 가중치가 w1=1, w2=1, 편향 b=-1.5이다.\n입력이 (1, 0)일 때, 가중합과 출력을 구하시오.\n(활성화: 0 초과이면 1, 이하이면 0)",
    "가중합 = 1x1 + 0x1 + (-1.5) = -0.5\n-0.5 <= 0 이므로 출력 = 0",
    img_name="perceptron_diagram.png"
)

# B-2
create_calculation_slide(
    "B-2", "MA", "3차시",
    "첫 번째 층이 y = 2x + 1이고, 두 번째 층이 z = 3y + 4일 때,\nz를 x에 대한 식으로 나타내시오.\n이 결과가 의미하는 바를 한 문장으로 쓰시오.",
    "z = 3(2x + 1) + 4 = 6x + 3 + 4 = 6x + 7 (여전히 1차식)\n의미: 활성화함수 없이 선형 층을 쌓으면,\n아무리 많이 쌓아도 하나의 선형 함수와 동일하다.",
    img_name="linear_composition.png"
)

# B-3
create_calculation_slide(
    "B-3", "MA", "4차시",
    "날씨 예측 AI의 3일간 예측: [4, 2, 6]\n실제 기온: [3, 2, 9]\nMSE(평균제곱오차)를 구하시오.",
    "오차: (4-3)=1, (2-2)=0, (6-9)=-3\n오차 제곱: 1, 0, 9\n합: 1+0+9 = 10\nMSE = 10/3 = 3.33",
    img_name="mse_outlier_comparison.png"
)

# B-4
create_calculation_slide(
    "B-4", "MA", "4차시",
    "데이터 A: 예측 [3,5,7], 실제 [2,5,8]\n데이터 B: 예측 [3,5,7], 실제 [2,5,20] (이상치)\nA와 B 각각의 MSE를 구하고, 이상치가 MSE에 미치는 영향을 설명하시오.",
    "A: [(1)2+(0)2+(-1)2]/3 = 2/3 = 0.67\nB: [(1)2+(0)2+(-13)2]/3 = 170/3 = 56.67\n이상치 하나로 MSE가 약 85배 커졌다.\nMSE는 오차를 제곱하므로 큰 오차에 민감하다.",
    img_name="mse_outlier_comparison.png"
)

# B-5
create_calculation_slide(
    "B-5", "MA", "5차시",
    "현재 가중치 w = 3.0, 기울기 = 4, 학습률 = 0.05일 때,\n한 번 업데이트하면 새 가중치는?",
    "w_new = 3.0 - 0.05 x 4 = 3.0 - 0.2 = 2.8",
    img_name="gradient_descent_comparison.png"
)

# B-6
create_calculation_slide(
    "B-6", "MA", "5차시",
    "w = 5.0, 기울기 = -3일 때,\n학습률 0.1 / 0.5 / 2.0 각각의 새 가중치를 구하고,\n어느 학습률이 가장 적절한지 이유와 함께 쓰시오.",
    "lr=0.1: 5.0 - 0.1x(-3) = 5.3 (안정적)\nlr=0.5: 5.0 - 0.5x(-3) = 6.5 (큰 이동)\nlr=2.0: 5.0 - 2.0x(-3) = 11.0 (발산 위험)\n적절: 0.1 또는 0.5. 2.0은 최솟값을 뛰어넘어 발산 가능.",
    img_name="gradient_descent_comparison.png"
)

# B-7
create_calculation_slide(
    "B-7", "MA", "6차시",
    "입력 x1=2, x2=1, 가중치 w1=0.5, w2=-1,\n편향 b=0.5, 활성화함수 ReLU일 때 출력을 구하시오.",
    "가중합: z = (2x0.5) + (1x-1) + 0.5 = 1 - 1 + 0.5 = 0.5\nReLU(0.5) = 0.5 (양수이므로 그대로)",
    img_name="forward_pass_example.jpg"
)

# B-8
create_calculation_slide(
    "B-8", "MA", "2~3차시",
    "다음 신경망의 학습 가능한 파라미터\n(가중치 + 편향)의 총 개수를 구하시오.\n구조: 입력 3 -> 은닉 4 -> 출력 2",
    "입력->은닉: 가중치 3x4=12개, 편향 4개 = 소계 16개\n은닉->출력: 가중치 4x2=8개, 편향 2개 = 소계 10개\n총 26개",
    img_name="nn_3_4_2.png"
)

# B-9
create_calculation_slide(
    "B-9", "MA", "7차시",
    "서울=[3,1], 한국=[2,0], 일본=[2,1]일 때,\n\"서울 - 한국 + 일본\"을 계산하고,\n이 결과가 의미하는 바를 한 문장으로 쓰시오.",
    "[3-2+2, 1-0+1] = [3, 2]\n의미: 서울이 한국의 수도이듯,\n결과 벡터는 일본의 수도(도쿄)에 해당하는 위치를 가리킨다.",
    img_name="word_vectors.png"
)

# --- 7. PART C TITLE ---
create_part_title_slide("Part C", "시각 / 분석형", "5문항 | 16점")

# C-1
create_question_slide(
    "C-1", "CO", "3차시",
    "위 4개의 그래프를 알맞은 활성화함수 이름과 연결하시오.\n\n보기: Sigmoid / ReLU / Tanh / Leaky ReLU\n\n정답: (가)=ReLU, (나)=Tanh, (다)=Sigmoid, (라)=Leaky ReLU\n\n채점: 4개 모두 정확 2점, 3개 1점, 2개 이하 0점",
    img_name="activation_functions_unlabeled.png"
)

# C-2
create_question_slide(
    "C-2", "AL", "6차시",
    "다음을 역전파 과정의 올바른 순서로 배열하시오.\n\n[ 가중치 업데이트 / 순전파 수행 / 기울기 역방향 전파 / 오차 계산 ]\n\n정답:\n1. 순전파 수행\n2. 오차 계산\n3. 기울기 역방향 전파\n4. 가중치 업데이트",
    img_name="backpropagation_flow.png"
)

# C-3
create_question_slide(
    "C-3", "MA + CO", "3차시",
    "(1) 위 좌표평면에 ReLU 함수의 그래프를 직접 그리시오. (2점)\n\n(2) 10층 이상의 깊은 신경망에서, Sigmoid 대신 ReLU를\n사용하는 이유를 '기울기 소실' 관점에서 설명하시오. (2점)\n\n핵심 키워드: 기울기 소실, Sigmoid 기울기 0, ReLU 기울기 1, 깊은 신경망",
    img_name="relu_template.png"
)

# C-4
create_question_slide(
    "C-4", "CO + AL", "3차시",
    "위 그래프를 참고하여, 다층 신경망에서 활성화함수 없이\n층을 쌓으면 왜 의미가 없는지 설명하고,\n활성화함수가 부여하는 핵심 성질의 이름을 쓰시오.\n\n핵심 키워드: 선형 함수 합성 = 선형, 비선형성, 공간 변환\n정답 키워드: 비선형성 (non-linearity)",
    img_name="linear_composition.png"
)

# C-5
create_question_slide(
    "C-5", "CO + MA", "2차시",
    "(1) 위 좌표평면의 XOR 데이터를 확인하시오:\n(0,0)=빨강, (0,1)=파랑, (1,0)=파랑, (1,1)=빨강 (1점)\n\n(2) 직선 하나로 빨강과 파랑을 완전히 분리할 수 있는지\n판단하고, 그 이유를 쓰시오. (3점)",
    img_name="xor_coordinate.png"
)

# --- 8. PART D TITLE ---
create_part_title_slide("Part D", "서술형 -- 개별 주제", "6문항 x 4점 = 24점 | 핵심 키워드 사전 공개")

# D-1
create_essay_slide(
    "D-1", "CT", "1차시",
    "튜링 테스트의 판정 기준과 중국어 방 반박을 각각 설명하고,\n현재의 ChatGPT에 대해 자신의 의견을 근거와 함께 쓰시오.",
    "튜링 테스트, 중국어 방, 구별 불가 = 지능, 이해 vs 흉내",
    img_name="turing_test_illustration.jpg",
    img_name2="chinese_room.jpg"
)

# D-2
create_essay_slide(
    "D-2", "CO + CT", "2+3차시",
    "단층 퍼셉트론이 XOR 문제를 풀 수 없는 이유를 설명하고,\n활성화함수가 이 한계를 어떻게 극복하는지 서술하시오.",
    "직선 분리 불가, 대각선 배치, 비선형성, 공간 변환, 다층 신경망",
    img_name="xor_coordinate.png",
    img_name2="nn_layers_comparison.jpg"
)

# D-3
create_essay_slide(
    "D-3", "CO", "4차시",
    "손실함수의 정의와 신경망 학습에서의 역할을 설명하시오.",
    "예측값과 실제값의 차이, 수치화, 손실 줄이기 = 학습, MSE",
    img_name="loss_function_target.jpg"
)

# D-4
create_essay_slide(
    "D-4", "CO + CT", "5차시",
    "경사하강법의 원리를 설명하고,\n학습률이 너무 클 때와 너무 작을 때\n각각 어떤 문제가 생기는지 서술하시오.",
    "기울기 반대 방향, 가중치 업데이트, 과대 -> 발산, 과소 -> 느린 수렴",
    img_name="gradient_descent_mountain.jpg",
    img_name2="gradient_descent_comparison.png"
)

# D-5
create_essay_slide(
    "D-5", "AL", "6차시",
    "순전파와 역전파의 관계를 4단계로 설명하시오.\n\n1. 순전파 (입력 -> 출력)\n2. 오차 계산\n3. 기울기 역전파 (출력 -> 입력)\n4. 가중치 업데이트",
    "순전파(입력->출력), 오차 계산, 기울기 역전파(출력->입력), 가중치 업데이트",
    img_name="backpropagation_flow.png"
)

# D-6
create_essay_slide(
    "D-6", "CO", "8차시",
    "LLM의 '다음 토큰 예측' 원리를 설명하고,\nTemperature가 출력에 미치는 영향을 서술하시오.",
    "확률 분포, 다음 토큰 예측, 낮으면 -> 반복적, 높으면 -> 다양",
    img_name="next_token_prediction.jpg",
    img_name2="temperature_comparison.jpg"
)

# --- 9. PART E TITLE ---
create_part_title_slide("Part E", "서술형 -- 통합", "2문항 x 5점 = 10점 | 차시 간 연결")

# E-1
create_essay_slide(
    "E-1", "AL", "4+5+6차시",
    "신경망이 하나의 데이터에서 학습하는 전체 과정을\n\"손실함수 -> 경사하강법 -> 역전파\"의 흐름으로\n연결하여 설명하시오.",
    "순전파->예측, 손실함수->오차 수치화, 경사하강법->기울기 반대 이동, 역전파->기울기 분배, 반복",
    img_name="backpropagation_flow.png"
)

# E-2
create_essay_slide(
    "E-2", "CT + CO", "7+8차시",
    "LLM이 환각(거짓 정보)을 생성하는 원인을\n토큰화, 확률적 다음 토큰 예측의 관점에서 설명하시오.\n\"AI가 멍청해서\"가 아닌, 구조적으로 왜 필연적인지 서술하시오.",
    "토큰 = 통계적 단위, 확률 기반 예측, \"그럴듯한\" != \"사실인\", 사실 검증 불가",
    img_name="ai_hallucination.jpg"
)

# --- 10. PART F ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
bg.fill.solid()
bg.fill.fore_color.rgb = PRIMARY_BLUE
bg.line.fill.background()

add_textbox(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(1),
            "Part F", font_size=60, bold=True, color=WHITE, font_name=FONT_NAME_EN)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.5), Inches(2.5), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = WHITE
line.line.fill.background()

add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.7),
            "자기 성찰 (Google Forms)", font_size=28, bold=True, color=WHITE)

# F-1 description
f_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(3.8), Inches(5), Inches(3)
)
f_box.fill.solid()
f_box.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x8A)
f_box.line.fill.background()
tf = f_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.2)
set_text(tf, "F-1. Before (1차시 시작)", font_size=18, bold=True, color=WHITE)
add_paragraph(tf, "\"기계가 생각할 수 있는가?\"", font_size=14, color=RGBColor(0xBF, 0xDB, 0xFE))
add_paragraph(tf, "1. 이 질문에 대한 나의 현재 답", font_size=13, color=WHITE)
add_paragraph(tf, "2. 그렇게 생각하는 근거", font_size=13, color=WHITE)
add_paragraph(tf, "3. AI에 대해 내가 이미 알고 있는 것", font_size=13, color=WHITE)
add_paragraph(tf, "4. 이 수업에서 가장 알고 싶은 것", font_size=13, color=WHITE)

# F-2 description
f_box2 = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(3.8), Inches(5), Inches(3)
)
f_box2.fill.solid()
f_box2.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x8A)
f_box2.line.fill.background()
tf2 = f_box2.text_frame
tf2.word_wrap = True
tf2.margin_left = Inches(0.3)
tf2.margin_top = Inches(0.2)
set_text(tf2, "F-2. After (10차시 종료)", font_size=18, bold=True, color=WHITE)
add_paragraph(tf2, "\"기계가 생각할 수 있는가?\" - 지금의 답", font_size=14, color=RGBColor(0xBF, 0xDB, 0xFE))
add_paragraph(tf2, "1. 지금의 답과 근거", font_size=13, color=WHITE)
add_paragraph(tf2, "2. 무엇이 변했나요?", font_size=13, color=WHITE)
add_paragraph(tf2, "3. 결정적이었던 수업 경험", font_size=13, color=WHITE)
add_paragraph(tf2, "4. 배운 개념 3개 이상 활용한 근거", font_size=13, color=WHITE)
add_paragraph(tf2, "5. AI와 인간의 미래에 대한 생각", font_size=13, color=WHITE)


# --- 11. ANSWER KEY SLIDES ---

# Answer Key: Part A
create_part_title_slide("Answer Key", "Part A 정답", "빈칸/단답 정답 및 허용 답안")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_top_bar(slide, color=GREEN)
add_bottom_bar(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
            "Part A 정답표", font_size=24, bold=True, color=GREEN)

answers_a = [
    ("A-1", "앨런 튜링", "Alan Turing, 튜링"),
    ("A-2", "퍼셉트론", "Perceptron"),
    ("A-3", "시그모이드 (Sigmoid)", "시그모이드 함수"),
    ("A-4", "렐루 (ReLU)", "렐루 함수"),
    ("A-5", "손실함수", "loss function, 비용함수"),
    ("A-6", "학습률", "learning rate"),
    ("A-7", "조정(변경), 유지(그대로)", "수정/업데이트, 놔둔다"),
    ("A-8", "학습률, 기울기(gradient)", "lr, 경사, gradient"),
    ("A-9", "BPE", "Byte Pair Encoding"),
    ("A-10", "환각", "Hallucination, 할루시네이션"),
]

# Table-like layout
for i, (qnum, answer, allowed) in enumerate(answers_a):
    y = Inches(1.2) + Inches(0.55) * i
    # Alternating row background
    if i % 2 == 0:
        row_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), y - Inches(0.05), Inches(12), Inches(0.50)
        )
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = VERY_LIGHT_BLUE
        row_bg.line.fill.background()

    add_textbox(slide, Inches(0.7), y, Inches(1), Inches(0.4),
                qnum, font_size=13, bold=True, color=DARK_BLUE, font_name=FONT_NAME_EN)
    add_textbox(slide, Inches(1.8), y, Inches(4), Inches(0.4),
                answer, font_size=13, bold=True, color=BLACK)
    add_textbox(slide, Inches(6.5), y, Inches(5.5), Inches(0.4),
                f"허용: {allowed}", font_size=12, color=GRAY)


# Answer Key: Part B
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_top_bar(slide, color=GREEN)
add_bottom_bar(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
            "Part B 정답 요약", font_size=24, bold=True, color=GREEN)

answers_b = [
    ("B-1", "가중합 = -0.5, 출력 = 0"),
    ("B-2", "z = 6x + 7 (여전히 1차식)"),
    ("B-3", "MSE = 10/3 = 3.33"),
    ("B-4", "A: 0.67 / B: 56.67 (85배 차이)"),
    ("B-5", "w_new = 2.8"),
    ("B-6", "0.1: 5.3 / 0.5: 6.5 / 2.0: 11.0"),
    ("B-7", "z = 0.5, ReLU = 0.5"),
    ("B-8", "총 26개 (16 + 10)"),
    ("B-9", "[3, 2] (도쿄에 해당)"),
]

for i, (qnum, answer) in enumerate(answers_b):
    y = Inches(1.2) + Inches(0.6) * i
    if i % 2 == 0:
        row_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), y - Inches(0.05), Inches(12), Inches(0.55)
        )
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = VERY_LIGHT_BLUE
        row_bg.line.fill.background()

    add_textbox(slide, Inches(0.7), y, Inches(1), Inches(0.45),
                qnum, font_size=14, bold=True, color=ORANGE, font_name=FONT_NAME_EN)
    add_textbox(slide, Inches(1.8), y, Inches(10), Inches(0.45),
                answer, font_size=14, color=BLACK)


# --- 12. RUBRIC SLIDES ---
create_part_title_slide("Rubrics", "서술형 채점 기준", "D, E 파트 5단계 루브릭")

# D rubric slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_top_bar(slide, color=DARK_BLUE)
add_bottom_bar(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
            "Part D 루브릭 요약", font_size=24, bold=True, color=DARK_BLUE)

rubric_d = [
    ("D-1", "튜링 테스트 vs 중국어 방", "양쪽 모두 정확 설명 + ChatGPT에 근거 2개 이상 자기 입장"),
    ("D-2", "XOR 한계와 활성화함수", "XOR 직선 분리 불가를 좌표 기반 설명 + 비선형 공간 변환 원리 연결"),
    ("D-3", "손실함수", "정의(예측-실제 차이 수치화) + 역할(학습 방향 제시) + 구체적 예시"),
    ("D-4", "경사하강법과 학습률", "원리 설명 + 과대(발산), 과소(느린 수렴) 문제 모두 서술"),
    ("D-5", "순전파와 역전파", "4단계 올바른 순서 + 연쇄법칙 역할까지 언급"),
    ("D-6", "다음 토큰 예측과 Temperature", "확률 분포 + 다음 토큰 예측 원리 + Temperature 영향"),
]

for i, (qnum, title, crit5) in enumerate(rubric_d):
    y = Inches(1.2) + Inches(0.95) * i
    if i % 2 == 0:
        row_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), y - Inches(0.05), Inches(12), Inches(0.90)
        )
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = VERY_LIGHT_BLUE
        row_bg.line.fill.background()

    add_textbox(slide, Inches(0.7), y, Inches(1), Inches(0.4),
                qnum, font_size=14, bold=True, color=DARK_BLUE, font_name=FONT_NAME_EN)
    add_textbox(slide, Inches(1.8), y, Inches(3), Inches(0.4),
                title, font_size=14, bold=True, color=BLACK)
    add_textbox(slide, Inches(1.8), y + Inches(0.4), Inches(10), Inches(0.4),
                f"5점: {crit5}", font_size=12, color=GRAY)


# E rubric slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_top_bar(slide, color=DARK_BLUE)
add_bottom_bar(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
            "Part E 루브릭 요약", font_size=24, bold=True, color=DARK_BLUE)

rubric_e = [
    ("E-1", "신경망 학습의 전체 과정 (4+5+6차시 통합)",
     "3개 개념을 인과적으로 연결하여 전체 학습 사이클을 완결적으로 설명 + \"반복\"까지",
     "2개 개념만 연결하거나, 3개 모두 언급하나 나열식",
     "1~2개 개념만 단편적으로 언급"),
    ("E-2", "환각의 구조적 원인 (7+8차시 통합)",
     "토큰화 + 확률적 예측 + \"구조적 필연\" 논리가 인과적으로 연결",
     "한 요소만 정확 또는 두 요소 모두 피상적",
     "\"확률로 맞추니까 틀릴 수 있다\" 수준"),
]

for i, (qnum, title, crit5, crit3, crit2) in enumerate(rubric_e):
    y = Inches(1.3) + Inches(2.8) * i

    # Question header
    add_textbox(slide, Inches(0.7), y, Inches(1), Inches(0.4),
                qnum, font_size=18, bold=True, color=DARK_BLUE, font_name=FONT_NAME_EN)
    add_textbox(slide, Inches(1.8), y, Inches(10), Inches(0.4),
                title, font_size=16, bold=True, color=BLACK)

    # Rubric levels
    levels = [
        ("5점", crit5, GREEN),
        ("3점", crit3, ORANGE),
        ("2점", crit2, RED),
    ]
    for j, (score, desc, color) in enumerate(levels):
        ly = y + Inches(0.6) + Inches(0.55) * j
        score_badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), ly, Inches(0.8), Inches(0.4)
        )
        score_badge.fill.solid()
        score_badge.fill.fore_color.rgb = color
        score_badge.line.fill.background()
        tf = score_badge.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = score
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME

        add_textbox(slide, Inches(2.8), ly, Inches(9), Inches(0.4),
                    desc, font_size=13, color=DARK_GRAY)


# --- 13. SCORING PRINCIPLES SLIDE ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_top_bar(slide)
add_bottom_bar(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
            "채점 원칙", font_size=28, bold=True, color=DARK_BLUE)

principles = [
    ("빈칸/단답", "정답 또는 허용 답안 일치 시 만점. 중간 점수 없음."),
    ("계산형", "풀이 과정 60% + 최종 답 40%. 과정이 맞으면 산술 실수 시 부분 점수(carried error)."),
    ("서술형", "3축 평가 -- (1) 핵심 키워드 포함(KW) (2) 논리적 설명(LG) (3) 적용/확장(AP)"),
    ("키워드 공개", "모든 서술형 문항의 핵심 키워드는 학생에게 미리 공개됨. 키워드가 빠지면 감점."),
    ("자기 성찰", "성실도(2) + 변화 인식(3) + 개념 활용(3) + 진정성(2) = 10점"),
]

for i, (title, desc) in enumerate(principles):
    y = Inches(1.3) + Inches(1.1) * i

    # Accent dot
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.1), Inches(0.15), Inches(0.15)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = PRIMARY_BLUE
    dot.line.fill.background()

    add_textbox(slide, Inches(1.4), y, Inches(3), Inches(0.4),
                title, font_size=16, bold=True, color=DARK_BLUE)
    add_textbox(slide, Inches(1.4), y + Inches(0.4), Inches(10.5), Inches(0.5),
                desc, font_size=14, color=DARK_GRAY)


# --- 14. FINAL SLIDE ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BLUE
bg.line.fill.background()

add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
            "AI 기초 과정중심평가 문항집", font_size=36, bold=True, color=WHITE)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.2), Inches(2), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = PRIMARY_BLUE
line.line.fill.background()

add_textbox(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.6),
            "고등학교 2학년  |  인공지능 기초  |  2026", font_size=18, color=RGBColor(0xBF, 0xDB, 0xFE))

add_textbox(slide, Inches(1.5), Inches(5.3), Inches(10), Inches(0.6),
            "총 34문항  |  A/B/C 버전  |  4대 테마(CT, CO, AL, MA)", font_size=16, color=RGBColor(0x93, 0xC5, 0xFD))


# === SAVE ===
prs.save(OUT_PATH)
print(f"\nDone! Saved to: {OUT_PATH}")
print(f"Total slides: {len(prs.slides)}")
